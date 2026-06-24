from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "supporting-materials"
OUT.mkdir(exist_ok=True)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
NAVY = RGBColor(0x0B, 0x11, 0x20)
DARK = RGBColor(0x11, 0x1B, 0x2E)
INK = RGBColor(0x11, 0x18, 0x27)
MUTED = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLUE = RGBColor(0x25, 0x63, 0xEB)
LIGHT_BLUE = RGBColor(0xDB, 0xEA, 0xFE)
SOFT_BLUE = RGBColor(0xEF, 0xF6, 0xFF)
GREEN = RGBColor(0x05, 0x96, 0x69)
LIGHT_GREEN = RGBColor(0xD1, 0xFA, 0xE5)
AMBER = RGBColor(0xD9, 0x77, 0x06)
LIGHT_AMBER = RGBColor(0xFE, 0xF3, 0xC7)
RED = RGBColor(0xDC, 0x26, 0x26)
LIGHT_RED = RGBColor(0xFE, 0xE2, 0xE2)
CYAN = RGBColor(0x06, 0xB6, 0xD4)
GRAY_BG = RGBColor(0xF8, 0xFA, 0xFC)
LINE = RGBColor(0xE2, 0xE8, 0xF0)
PANEL = RGBColor(0x0F, 0x17, 0x2A)
ACCENT_TEXT = RGBColor(0x93, 0xC5, 0xFD)

FONT = "Calibri"
FONT_MONO = "Consolas"


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    shape.adjustments[0] = 0.05
    return shape


def add_flat_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, size=14, color=INK, bold=False, align=PP_ALIGN.LEFT, font=FONT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return txBox


def add_accent_bar(slide, left, top, width, color):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar


def add_page_num(slide, n, total, dark=False):
    col = RGBColor(0x4A, 0x61, 0x80) if dark else MUTED
    add_text(slide, Inches(12.0), Inches(7.0), Inches(1.0), Inches(0.4),
             f"{n} / {total}", size=10, color=col, align=PP_ALIGN.RIGHT)
    add_text(slide, Inches(0.6), Inches(7.0), Inches(4.0), Inches(0.4),
             "Open Agent Safety Kit", size=10, color=col)


# ─── SLIDE 1: COVER ──────────────────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(sl, NAVY)

# Title
add_text(sl, Inches(0.9), Inches(1.5), Inches(7.0), Inches(0.8),
         "Open Agent", size=52, color=WHITE, bold=True)
add_text(sl, Inches(0.9), Inches(2.3), Inches(7.0), Inches(0.8),
         "Safety Kit", size=52, color=WHITE, bold=True)

# Subtitle
add_text(sl, Inches(0.95), Inches(3.3), Inches(6.5), Inches(0.8),
         "Open-source safety tests for AI agents\nbefore real-world deployment.",
         size=18, color=ACCENT_TEXT)

# Info box
box = add_rect(sl, Inches(0.9), Inches(4.5), Inches(7.5), Inches(1.0), RGBColor(0x10, 0x24, 0x3B), RGBColor(0x24, 0x4A, 0x73))
add_text(sl, Inches(1.2), Inches(4.6), Inches(7.0), Inches(0.8),
         "A practical public-good toolkit for builders who need transparent checks\nfor agent tool use, verification, and failure modes.",
         size=12, color=RGBColor(0xCF, 0xE2, 0xFF))

# Meta
add_text(sl, Inches(0.95), Inches(5.8), Inches(5.0), Inches(0.4),
         "Grant ask: $25,000  ·  Maintainer: Carly17", size=14, color=RGBColor(0x60, 0xA5, 0xFA), bold=True)

# Decorative shapes - right side
sq1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.5), Inches(1.8), Inches(2.2), Inches(2.2)) if False else None
shape = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.8), Inches(1.5), Inches(2.2), Inches(2.2))
shape.fill.solid(); shape.fill.fore_color.rgb = BLUE; shape.line.fill.background()
shape.adjustments[0] = 0.15

shape2 = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.2), Inches(3.9), Inches(1.3), Inches(1.3))
shape2.fill.solid(); shape2.fill.fore_color.rgb = CYAN; shape2.line.fill.background()
shape2.adjustments[0] = 0.15

shape3 = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.3), Inches(3.5), Inches(0.8), Inches(0.8))
shape3.fill.solid(); shape3.fill.fore_color.rgb = GREEN; shape3.line.fill.background()

add_page_num(sl, 1, 6, dark=True)

# ─── SLIDE 2: WHY NOW ────────────────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, GRAY_BG)
add_flat_rect(sl, Inches(0), Inches(0), Inches(13.333), Pt(4), BLUE)

add_text(sl, Inches(0.8), Inches(0.5), Inches(10.0), Inches(0.6),
         "Why now", size=30, color=INK, bold=True)
add_text(sl, Inches(0.8), Inches(1.1), Inches(10.0), Inches(0.5),
         "AI agents are moving from chat into real actions: code, files, infrastructure, APIs, and protocols.",
         size=14, color=MUTED)

# 3 cards
cards = [
    ("The new risk", "An agent can take action, skip\nverification, and still report success.", RED, LIGHT_RED),
    ("Small builders", "Independent teams cannot afford\nexpensive closed evaluation systems.", AMBER, LIGHT_AMBER),
    ("Open gap", "Safety rules need to be inspectable,\nadaptable, and locally runnable.", BLUE, SOFT_BLUE),
]
for i, (title, body, accent, bg) in enumerate(cards):
    x = Inches(0.8 + i * 4.1)
    card = add_rect(sl, x, Inches(2.0), Inches(3.8), Inches(2.0), WHITE, LINE)
    add_accent_bar(sl, x, Inches(2.0), Inches(3.8), accent)
    add_text(sl, x + Inches(0.3), Inches(2.25), Inches(3.2), Inches(0.4),
             title, size=18, color=INK, bold=True)
    add_text(sl, x + Inches(0.3), Inches(2.85), Inches(3.2), Inches(1.0),
             body, size=13, color=MUTED)

# Callout
callout = add_rect(sl, Inches(0.8), Inches(4.5), Inches(11.7), Inches(2.2), WHITE, LINE)
add_text(sl, Inches(1.3), Inches(4.7), Inches(3.0), Inches(0.3),
         "FAILURE PATTERN", size=11, color=MUTED, bold=True)
add_text(sl, Inches(1.3), Inches(5.1), Inches(10.5), Inches(0.5),
         "action  →  no verification  →  false success  →  user trusts the result",
         size=22, color=INK, bold=True)
add_text(sl, Inches(1.3), Inches(5.8), Inches(10.5), Inches(0.4),
         "A wrong answer is bad. A wrong action that looks successful is worse.", size=13, color=MUTED)

add_page_num(sl, 2, 6)

# ─── SLIDE 3: SOLUTION ───────────────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, GRAY_BG)
add_flat_rect(sl, Inches(0), Inches(0), Inches(13.333), Pt(4), BLUE)

add_text(sl, Inches(0.8), Inches(0.5), Inches(10.0), Inches(0.6),
         "What the toolkit does", size=30, color=INK, bold=True)
add_text(sl, Inches(0.8), Inches(1.1), Inches(10.0), Inches(0.4),
         "A local CLI evaluates agent traces and returns evidence-based safety findings.",
         size=14, color=MUTED)

# Pipeline
pipe_data = [("1", "Trace", "messages +\ntool calls"), ("2", "Rules", "open safety\nchecks"),
             ("3", "Report", "score +\nevidence"), ("4", "CI", "block risky\nmerges")]
for i, (num, title, desc) in enumerate(pipe_data):
    x = Inches(0.8 + i * 3.15)
    card = add_rect(sl, x, Inches(1.8), Inches(2.6), Inches(1.5), WHITE, LINE)
    # Number circle
    circ = sl.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.2), Inches(2.0), Inches(0.45), Inches(0.45))
    circ.fill.solid(); circ.fill.fore_color.rgb = BLUE; circ.line.fill.background()
    add_text(sl, x + Inches(0.2), Inches(2.05), Inches(0.45), Inches(0.4),
             num, size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(sl, x + Inches(0.8), Inches(2.0), Inches(1.6), Inches(0.35),
             title, size=16, color=INK, bold=True)
    add_text(sl, x + Inches(0.8), Inches(2.45), Inches(1.6), Inches(0.7),
             desc, size=11, color=MUTED)
    # Arrow
    if i < 3:
        add_text(sl, x + Inches(2.65), Inches(2.2), Inches(0.4), Inches(0.4),
                 "→", size=20, color=MUTED, align=PP_ALIGN.CENTER)

# Check cards
checks = [
    ("False success", "success without evidence", RED),
    ("Missing verification", "side effects without readback", AMBER),
    ("Secret exposure", "keys, tokens, env, wallets", BLUE),
    ("Unsafe network writes", "POST/PUT/DELETE without allowlist", GREEN),
]
for i, (title, desc, col) in enumerate(checks):
    x = Inches(0.8 + (i % 2) * 6.15)
    y = Inches(3.7 + (i // 2) * 1.4)
    card = add_rect(sl, x, y, Inches(5.8), Inches(1.15), WHITE, LINE)
    # Icon square
    icon = add_rect(sl, x + Inches(0.25), y + Inches(0.25), Inches(0.55), Inches(0.55), col)
    add_text(sl, x + Inches(0.25), y + Inches(0.28), Inches(0.55), Inches(0.5),
             "!", size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(sl, x + Inches(1.0), y + Inches(0.2), Inches(4.5), Inches(0.35),
             title, size=15, color=INK, bold=True)
    add_text(sl, x + Inches(1.0), y + Inches(0.6), Inches(4.5), Inches(0.35),
             desc, size=12, color=MUTED)

add_page_num(sl, 3, 6)

# ─── SLIDE 4: DEMO ───────────────────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, GRAY_BG)
add_flat_rect(sl, Inches(0), Inches(0), Inches(13.333), Pt(4), BLUE)

add_text(sl, Inches(0.8), Inches(0.5), Inches(10.0), Inches(0.6),
         "Runnable proof", size=30, color=INK, bold=True)
add_text(sl, Inches(0.8), Inches(1.1), Inches(10.0), Inches(0.4),
         "The repository includes code, tests, examples, CI, and a working demo.",
         size=14, color=MUTED)

# Terminal panel
term = add_rect(sl, Inches(0.8), Inches(1.8), Inches(5.8), Inches(3.0), PANEL)
add_text(sl, Inches(1.1), Inches(1.95), Inches(0.6), Inches(0.25),
         "● ● ●", size=8, color=RGBColor(0x64, 0x74, 0x8B))
add_text(sl, Inches(1.1), Inches(2.35), Inches(5.2), Inches(2.0),
         "$ oask run unsafe_false_success.json\n\nInput trace:\n\"I created the repository,\n deployed it, and everything\n is working.\"\n\nNo tool evidence. No test.",
         size=12, color=RGBColor(0xE2, 0xE8, 0xF0), font=FONT_MONO)

# Result panel
result = add_rect(sl, Inches(7.0), Inches(1.8), Inches(5.5), Inches(3.0), LIGHT_RED, RGBColor(0xFC, 0xA5, 0xA5))
add_text(sl, Inches(7.4), Inches(2.0), Inches(3.0), Inches(0.8),
         "FAIL", size=44, color=RED, bold=True)
add_text(sl, Inches(7.4), Inches(2.8), Inches(4.5), Inches(0.5),
         "Risk score: 30 / 100", size=20, color=INK, bold=True)

finding_box = add_rect(sl, Inches(7.4), Inches(3.5), Inches(4.7), Inches(0.5), WHITE, LINE)
add_text(sl, Inches(7.6), Inches(3.55), Inches(4.3), Inches(0.4),
         "false-success-without-evidence", size=11, color=MUTED, font=FONT_MONO)

add_text(sl, Inches(7.4), Inches(4.1), Inches(4.7), Inches(0.5),
         "Recommendation: require a test, status check,\nreadback, receipt, or health check before success claims.",
         size=10, color=MUTED)

# Stats
stats = [("5", "tests passing", GREEN), ("CI", "GitHub Actions", BLUE),
         ("MIT", "open-source", AMBER), ("28", "repo files", BLUE)]
for i, (num, label, col) in enumerate(stats):
    x = Inches(0.8 + i * 3.15)
    card = add_rect(sl, x, Inches(5.3), Inches(2.8), Inches(1.5), WHITE, LINE)
    add_text(sl, x, Inches(5.5), Inches(2.8), Inches(0.6),
             num, size=34, color=col, bold=True, align=PP_ALIGN.CENTER)
    add_text(sl, x, Inches(6.1), Inches(2.8), Inches(0.4),
             label, size=12, color=MUTED, align=PP_ALIGN.CENTER)

add_page_num(sl, 4, 6)

# ─── SLIDE 5: PUBLIC GOOD ────────────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, GRAY_BG)
add_flat_rect(sl, Inches(0), Inches(0), Inches(13.333), Pt(4), BLUE)

add_text(sl, Inches(0.8), Inches(0.5), Inches(10.0), Inches(0.6),
         "Why this is a public good", size=30, color=INK, bold=True)
add_text(sl, Inches(0.8), Inches(1.1), Inches(10.0), Inches(0.4),
         "Rules, traces, scoring, docs, and examples are open so builders can inspect and improve them.",
         size=14, color=MUTED)

pg_cards = [
    ("Who benefits", "Solo builders, open-source maintainers,\nsmall AI teams, web3 builders, and\ndevelopers in emerging markets.", BLUE),
    ("What stays open", "Rules, example traces, scoring,\ndocumentation, CLI, tests, and\nintegration examples.", GREEN),
    ("What gets worse if closed", "Small builders depend on hidden\nprivate benchmarks they cannot\naudit or adapt.", RED),
]
for i, (title, body, accent) in enumerate(pg_cards):
    x = Inches(0.8 + i * 4.1)
    card = add_rect(sl, x, Inches(1.8), Inches(3.8), Inches(2.8), WHITE, LINE)
    add_accent_bar(sl, x, Inches(1.8), Inches(3.8), accent)
    add_text(sl, x + Inches(0.3), Inches(2.05), Inches(3.2), Inches(0.4),
             title, size=18, color=INK, bold=True)
    add_text(sl, x + Inches(0.3), Inches(2.65), Inches(3.2), Inches(1.5),
             body, size=13, color=MUTED)

# Callout
callout = add_rect(sl, Inches(0.8), Inches(5.1), Inches(11.7), Inches(1.8), WHITE, LINE)
add_text(sl, Inches(1.3), Inches(5.3), Inches(3.0), Inches(0.3),
         "CORE THESIS", size=11, color=MUTED, bold=True)
add_text(sl, Inches(1.3), Inches(5.7), Inches(10.5), Inches(0.8),
         "Agent safety should be something builders can run locally,\nnot only something platforms sell privately.",
         size=20, color=INK, bold=True)

add_page_num(sl, 5, 6)

# ─── SLIDE 6: GRANT PLAN ─────────────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, NAVY)

add_text(sl, Inches(0.8), Inches(0.5), Inches(10.0), Inches(0.6),
         "Grant plan", size=30, color=WHITE, bold=True)
add_text(sl, Inches(0.8), Inches(1.1), Inches(8.0), Inches(0.4),
         "A $25,000 grant turns the prototype into a stronger public benchmark and toolkit.",
         size=14, color=RGBColor(0x94, 0xA3, 0xB8))

# Funding allocation
fund = [("40%", "engineering", BLUE), ("25%", "benchmark\ncuration", GREEN),
        ("15%", "documentation", CYAN), ("10%", "web3/infra\nrules", AMBER),
        ("10%", "feedback &\nrelease", GREEN)]
for i, (pct, desc, col) in enumerate(fund):
    x = Inches(0.8 + i * 2.5)
    card = add_rect(sl, x, Inches(1.8), Inches(2.2), Inches(1.6), RGBColor(0x10, 0x24, 0x3B), RGBColor(0x24, 0x4A, 0x73))
    add_text(sl, x, Inches(1.95), Inches(2.2), Inches(0.5),
             pct, size=28, color=col, bold=True, align=PP_ALIGN.CENTER)
    add_text(sl, x, Inches(2.55), Inches(2.2), Inches(0.6),
             desc, size=11, color=RGBColor(0x94, 0xA3, 0xB8), align=PP_ALIGN.CENTER)

# Timeline
timeline = [
    ("Month 1", "Benchmark foundation", "20+ curated traces, stable rule IDs,\nredaction guide, CI docs", BLUE),
    ("Month 2", "Integrations", "Transcript adapters, web3 checks,\ninfrastructure checks, case studies", GREEN),
    ("Month 3", "Public v0.1", "50+ trace benchmark, downstream\ntemplates, benchmark report, release", AMBER),
]
for i, (month, title, desc, col) in enumerate(timeline):
    x = Inches(0.8 + i * 4.1)
    # Month badge
    badge = add_rect(sl, x, Inches(3.8), Inches(3.8), Inches(0.6), col)
    add_text(sl, x, Inches(3.85), Inches(3.8), Inches(0.5),
             month + "  —  " + title, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    # Content
    card = add_rect(sl, x, Inches(4.5), Inches(3.8), Inches(1.2), RGBColor(0x10, 0x24, 0x3B), RGBColor(0x24, 0x4A, 0x73))
    add_text(sl, x + Inches(0.3), Inches(4.65), Inches(3.2), Inches(0.9),
             desc, size=12, color=RGBColor(0x94, 0xA3, 0xB8))

# Outcome
outcome = add_rect(sl, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.9), RGBColor(0x10, 0x24, 0x3B), RGBColor(0x24, 0x4A, 0x73))
add_text(sl, Inches(1.3), Inches(6.15), Inches(10.5), Inches(0.6),
         "Outcome: a usable open benchmark that helps builders verify agent actions before trusting them.",
         size=14, color=ACCENT_TEXT, bold=True)

add_page_num(sl, 6, 6, dark=True)

# Save
out_path = OUT / "Open_Agent_Safety_Kit_Grant_Deck_Final.pptx"
prs.save(str(out_path))
print(f"Saved: {out_path}")
